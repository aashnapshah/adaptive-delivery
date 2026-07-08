"""Add one screenshot slide per pipeline stage to research_plan/Sandbox.pptx.

Each slide = a title + the headless screenshot of one tab of ui/pipeline.html
(captured into ui/screenshots/). Re-running is idempotent: slides previously added
by this script (tagged in the title) are removed first, so you always get exactly
one fresh slide per stage.

SAFETY: refuses to run while PowerPoint has the deck open (a ``~$Sandbox.pptx`` lock
file is present), because PowerPoint would overwrite these edits on its next save.
Close the deck in PowerPoint, then:  python3 update_deck.py

Regenerate the screenshots first (optional) with:  python3 update_deck.py --shots
which re-renders every tab via headless Chrome before rebuilding the slides.
"""

from __future__ import annotations

import os
import subprocess
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.abspath(os.path.join(HERE, "..", "research_plan", "Sandbox.pptx"))
LOCK = os.path.join(os.path.dirname(DECK), "~$Sandbox.pptx")
SHOTS = os.path.join(HERE, "ui", "screenshots")
TAG = "[pipeline-demo]"   # marker placed in each generated slide's title text

# stage id -> (slide title, screenshot file)
STAGES = [
    ("Step 0 - Case picker (shared dataset)",            "0-case.png"),
    ("Stage 01 - Gatekeeper environment",                "1-env.png"),
    ("Stage 02 - Recommendation generation (3 configs)", "2-gen.png"),
    ("Stage 02 - Compare vs actually-ordered",           "3-cmp.png"),
    ("Stage 03 - Appropriateness & harm scoring",        "4-score.png"),
    ("Stage 04 - Delivery interface & response",         "5-iface.png"),
    ("Stage 05 - Adaptive delivery policy",              "6-policy.png"),
]

TABS = [("case", 900), ("env", 620), ("gen", 760), ("cmp", 700),
        ("score", 1080), ("iface", 1080), ("policy", 1000)]
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"


def shoot(port: int = 8731) -> None:
    """Re-render every tab to ui/screenshots/ via headless Chrome (server must be up)."""
    os.makedirs(SHOTS, exist_ok=True)
    for i, (tab, h) in enumerate(TABS):
        out = os.path.join(SHOTS, f"{i}-{tab}.png")
        url = f"http://localhost:{port}/ui/pipeline.html?case=ehr&tab={tab}"
        subprocess.run([CHROME, "--headless=new", "--disable-gpu", "--hide-scrollbars",
                        "--force-device-scale-factor=2", "--virtual-time-budget=4500",
                        f"--window-size=1480,{h}", f"--screenshot={out}", url],
                       check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        print("  rendered", out)


def _drop_existing(prs: Presentation) -> int:
    """Remove slides previously added by this script, so re-runs stay idempotent."""
    xml_slides = prs.slides._sldIdLst
    ids = list(xml_slides)
    removed = 0
    for i in range(len(prs.slides) - 1, -1, -1):
        slide = prs.slides[i]
        text = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        if TAG in text:
            xml_slides.remove(ids[i])
            removed += 1
    return removed


def build() -> None:
    if os.path.exists(LOCK):
        sys.exit(f"REFUSING: {LOCK} exists -> Sandbox.pptx is open in PowerPoint. "
                 "Close it, then re-run.")
    missing = [f for _, f in STAGES if not os.path.exists(os.path.join(SHOTS, f))]
    if missing:
        sys.exit(f"Missing screenshots: {missing}. Run with --shots (server up) first.")

    prs = Presentation(DECK)
    sw, sh = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]   # blank layout
    dropped = _drop_existing(prs)

    for title, fname in STAGES:
        slide = prs.slides.add_slide(blank)
        # title
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), sw - Inches(0.8), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run(); run.text = f"{title}  {TAG}"
        run.font.size = Pt(22); run.font.bold = True
        run.font.color.rgb = RGBColor(0x1B, 0x4F, 0x82)
        # image, scaled to fit under the title, centered
        img = os.path.join(SHOTS, fname)
        from PIL import Image  # noqa
        try:
            iw, ih = Image.open(img).size
        except Exception:
            iw, ih = 1480, 900
        top = Inches(0.95)
        avail_w = sw - Inches(0.8)
        avail_h = sh - top - Inches(0.3)
        scale = min(avail_w / iw, avail_h / ih)
        w, h = int(iw * scale), int(ih * scale)
        left = int((sw - w) / 2)
        slide.shapes.add_picture(img, left, top, width=w, height=h)

    prs.save(DECK)
    print(f"Removed {dropped} old demo slide(s); added {len(STAGES)}. Saved {DECK}")


if __name__ == "__main__":
    if "--shots" in sys.argv:
        shoot()
    build()
